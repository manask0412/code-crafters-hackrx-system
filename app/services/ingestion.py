# app/services/ingestion.py
import aiohttp
import tempfile
import os
import fitz  # PyMuPDF
import docx
import email
from email import policy
import tiktoken
import zipfile
import pytesseract
from PIL import Image
import pandas as pd
from pptx import Presentation
from urllib.parse import urlparse, unquote

from app.services.embeddings import upsert_chunks

# Configuration constants
ENCODING_MODEL = "text-embedding-ada-002"
CHUNK_TOKEN_SIZE = 2000
OVERLAP = int(0.1 * CHUNK_TOKEN_SIZE)
IMAGE_EXTS = {".png", ".jpeg", ".jpg", ".bmp", ".gif"}
ARCHIVE_EXTS = {".zip"}
EXCEL_EXTS = {".xlsx"}
PPTX_EXTS = {".pptx"}
TEXT_EXTS = {".pdf", ".docx", ".txt", ".eml"}

# Helpers for extension detection
async def get_url_extension(url: str) -> str:
    # 1) Try Content-Disposition header
    async with aiohttp.ClientSession() as session:
        async with session.get(url) as resp:
            resp.raise_for_status()
            cd = resp.headers.get("Content-Disposition", "")
            filename = ""
            if "filename=" in cd:
                filename = cd.split("filename=")[-1].strip("\"")
                return os.path.splitext(filename)[1].lower()
            else:
                # 2) Fallback to URL path
                path = urlparse(url).path
                return os.path.splitext(path)[1].lower()

async def download_to_temp(url: str, ext: str) -> str:
    timeout = aiohttp.ClientTimeout(total=60)  # Optional: limit total download time
    async with aiohttp.ClientSession(timeout=timeout) as session:
        async with session.get(url) as resp:
            resp.raise_for_status()
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as f:
                async for chunk in resp.content.iter_chunked(1024 * 1024):
                    f.write(chunk)
                return f.name

# Handlers
async def handle_bin_stub(basename: str, ext: str, doc_url: str):
    text = f"Binary file {basename}{ext}; content not extractable."
    await upsert_chunks([{"id": f"{basename}::metadata", "text": text}], doc_url)

async def handle_zip_stub(temp_path: str, basename: str, doc_url: str):
    with zipfile.ZipFile(temp_path) as z:
        entries = z.namelist()
    summary = f"Archive contains {len(entries)} files: {', '.join(entries)}"
    await upsert_chunks([{"id": f"{basename}::metadata", "text": summary}], doc_url)

async def handle_image_ocr(temp_path: str, basename: str, ext: str, doc_url: str):
    text = pytesseract.image_to_string(Image.open(temp_path)).strip()
    if len(text) >= 1:
        await _chunk_and_upsert(text, basename, doc_url)
    else:
        stub = f"Image file {basename}{ext}, OCR returned no useful text."
        await upsert_chunks([{"id": f"{basename}::metadata", "text": stub}], doc_url)

async def handle_xlsx(temp_path: str, basename: str, doc_url: str):
    book = pd.read_excel(temp_path, sheet_name=None, dtype=str)
    combined = []

    for sheet_name, df in book.items():
        df = df.fillna("").astype(str)
        row_texts = df.agg(" ".join, axis=1).tolist()
        sheet_text = f"Sheet {sheet_name}:\n" + "\n".join(row_texts)
        combined.append(sheet_text)

    all_text = "\n\n".join(combined)
    await _chunk_and_upsert(all_text, basename, doc_url)

async def handle_pptx(temp_path: str, basename: str, doc_url: str):
    prs = Presentation(temp_path)
    slides = []

    for idx, slide in enumerate(prs.slides, 1):
        text_runs = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
        slide_text = f"Slide {idx}: " + ' '.join(text_runs)
        slides.append(slide_text)

    full_text = "\n".join(slides)
    await _chunk_and_upsert(full_text, basename, doc_url)

# Text extraction
async def extract_text_file(path: str, ext: str) -> str:
    if ext == ".pdf":
        with fitz.open(path) as doc:
            return "\n".join(page.get_text() for page in doc)
    if ext == ".docx":
        doc = docx.Document(path)
        return "\n".join(p.text for p in doc.paragraphs)
    if ext == ".txt":
        return open(path, 'r', encoding='utf-8', errors='ignore').read()
    if ext == ".eml":
        with open(path, 'rb') as f:
            msg = email.message_from_binary_file(f, policy=policy.default)
        parts = []
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    parts.append(part.get_content())
        else:
            parts.append(msg.get_content())
        return "\n".join(parts)
    return ""

async def _chunk_and_upsert(text: str, basename: str, doc_url: str):
    enc = tiktoken.encoding_for_model(ENCODING_MODEL)
    tokens = enc.encode(" ".join(text.split()))
    records = []
    for i in range(0, len(tokens), CHUNK_TOKEN_SIZE - OVERLAP):
        chunk = tokens[i:i+CHUNK_TOKEN_SIZE]
        records.append({"id": f"{basename}::chunk_{i}", "text": enc.decode(chunk)})
    await upsert_chunks(records, doc_url)

# Main pipeline
async def extract_and_embed_chunks(doc_url: str) -> None:
    ext = await get_url_extension(doc_url)
    basename = unquote(os.path.splitext(os.path.basename(urlparse(doc_url).path))[0])

    # Binary or oversized
    if ext == ".bin":
        await handle_bin_stub(basename, ext, doc_url)
        return

    temp_path = await download_to_temp(doc_url, ext)
    try:
        if ext in ARCHIVE_EXTS:
            await handle_zip_stub(temp_path, basename, doc_url)
        elif ext in IMAGE_EXTS:
            await handle_image_ocr(temp_path, basename, ext, doc_url)
        elif ext in EXCEL_EXTS:
            await handle_xlsx(temp_path, basename, doc_url)
        elif ext in PPTX_EXTS:
            await handle_pptx(temp_path, basename, doc_url)
        elif ext in TEXT_EXTS:
            text = await extract_text_file(temp_path, ext)
            await _chunk_and_upsert(text, basename, doc_url)
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    finally:
        os.remove(temp_path)
